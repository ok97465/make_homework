r"""PPT 작성을 위한 Helper.

Todo:
    None

References:
    .. [] 책: 저자명. (년). 챕터명. In 편집자명 (역할), 책명 (쪽). 발행지 : 발행사
    .. [] 학위 논문: 학위자명, "논문제목", 대학원 이름 석사 학위논문, 1990
    .. [] 저널 논문: 저자. "논문제목". 저널명, . pp.

:auther: ok97465
:Date created: 22.08.15 16:03:12
"""
# %% Import
# Standard library imports
from typing import Literal

# Third party imports
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.dml.line import LineFormat
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN
from pptx.slide import Slide
from pptx.util import Cm, Pt

# %% Define Variable
OPERATOR_TYPE = Literal["+", "-", "×", "÷"]
GET_ANS_OPERATOR = {
    "+": lambda x, y: x + y,
    "-": lambda x, y: x - y,
    "×": lambda x, y: x * y,
    "÷": lambda x, y: x / y,
}


# %% Function
def get_new_ppt(width: float = 21.0, height: float = 29.7) -> Presentation:
    """Generate presentation instance."""
    prs = Presentation()

    prs.slide_width = Cm(width)
    prs.slide_height = Cm(height)

    return prs


def add_slide_title(slide: Slide, txt: str, slide_width: int):
    """Add slide title."""
    tb = slide.shapes.add_textbox(0, Cm(0.508), Cm(slide_width), Cm(1.905))

    tf = tb.text_frame

    tf.text = txt
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.size = Pt(40)
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, 0, Cm(2.5), Cm(slide_width), Cm(2.5)
    )


def add_vquiz(
    slide: Slide,
    v1: int,
    v2: int,
    add_answer: bool,
    left: float,
    top: float,
    operator: OPERATOR_TYPE,
    do_sort: bool = False,
):
    """Generate the vertical quiz for adding."""
    font_size = 30
    font_name = "Courier New"

    if do_sort:
        v1, v2 = sorted([v1, v2], reverse=True)

    tb = slide.shapes.add_textbox(Cm(left), Cm(top), Cm(left + 2.57), Cm(top + 3.9))
    tf = tb.text_frame

    tf.text = f"{v1:4d}\n{operator}{v2:3d}"

    if add_answer is True:
        tf.text += f"\n{GET_ANS_OPERATOR[operator](v1, v2):4d}"
        tf.paragraphs[2].font.size = Pt(font_size)
        tf.paragraphs[2].font.name = font_name

    tf.paragraphs[0].font.size = Pt(font_size)
    tf.paragraphs[0].font.name = font_name
    tf.paragraphs[1].font.size = Pt(font_size)
    tf.paragraphs[1].font.name = font_name

    tf.word_wrap = False
    tf.margin_top = Cm(0.0)
    tf.margin_bottom = Cm(0.0)
    tf.margin_left = Cm(0.0)
    tf.margin_right = Cm(0.0)

    l_left = left
    l_right = l_left + 2.57
    l_top = top + 2.7
    connector = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Cm(l_left),
        Cm(l_top),
        Cm(l_right),
        Cm(l_top),
    )

    # Black line
    if not hasattr(connector, "ln"):
        connector.ln = connector.get_or_add_ln()

    line = LineFormat(connector)
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(0, 0, 0)


def add_hquiz(
    slide: Slide,
    v1: int,
    v2: int,
    add_answer: bool,
    left: float,
    top: float,
    operator: OPERATOR_TYPE,
    do_sort: bool = False,
):
    """Generate the horizontal quiz for adding."""
    font_size = 30
    font_name = "Courier New"

    if do_sort:
        v1, v2 = sorted([v1, v2], reverse=True)

    tb = slide.shapes.add_textbox(Cm(left), Cm(top), Cm(left + 4), Cm(top + 1.3))
    tf = tb.text_frame

    tf.text = f"☞{v1:2d} {operator}{v2:2d} = "

    if add_answer is True:
        tf.text += f"{GET_ANS_OPERATOR[operator](v1, v2):2d}"

    tf.paragraphs[0].font.size = Pt(font_size)
    tf.paragraphs[0].font.name = font_name

    tf.word_wrap = False
    tf.margin_top = Cm(0.0)
    tf.margin_bottom = Cm(0.0)
    tf.margin_left = Cm(0.0)
    tf.margin_right = Cm(0.0)
